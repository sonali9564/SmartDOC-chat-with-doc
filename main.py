import os
import concurrent.futures
import time
from dotenv import load_dotenv
import streamlit as st
from langchain_community.document_loaders import UnstructuredPDFLoader
from langchain_text_splitters.character import CharacterTextSplitter
from langchain_community.vectorstores import FAISS
from langchain_community.embeddings import HuggingFaceEmbeddings
from langchain_groq import ChatGroq
from langchain.chains import ConversationalRetrievalChain
from docx import Document  # For DOCX files
import pythoncom
from pptx import Presentation  # For PowerPoint files
import re
from nltk.corpus import stopwords
import fitz  # PyMuPDF for better PDF extraction
import io
from streamlit_webrtc import webrtc_streamer, AudioProcessorBase
import speech_recognition as sr


# Load environment variables
load_dotenv()

# Define the working directory based on the current script location
working_dir = os.path.dirname(os.path.abspath(__file__))


# Define CustomDocument class
class CustomDocument:
    def __init__(self, page_content, metadata=None):
        self.page_content = page_content
        self.metadata = metadata or {}


# Function to load documents based on file type (PDF, Word, PowerPoint, Text)
def load_document(file_path):
    file_extension = file_path.split('.')[-1].lower()

    # Process file based on its type
    if file_extension == 'pdf':
        raw_documents = load_pdf(file_path)
    elif file_extension == 'docx':
        raw_documents = load_word(file_path)
    elif file_extension == 'pptx':
        raw_documents = load_ppt(file_path)
    elif file_extension == 'txt':
        raw_documents = load_txt(file_path)
    else:
        raise ValueError(f"Unsupported file type: {file_extension}")

    # Convert raw documents to CustomDocument objects
    documents = [CustomDocument(text) for text in raw_documents]
    return documents


# Function to load text from PDF using PyMuPDF (fitz) with more efficient extraction
def load_pdf(file_path):
    doc = fitz.open(file_path)
    text = []
    for page in doc:
        # Extract text more efficiently from each page
        page_text = page.get_text("text")
        text.append(page_text)
    return text


# Function to load text from Word (DOCX)
def load_word(file_path):
    doc = Document(file_path)
    text = []
    for paragraph in doc.paragraphs:
        text.append(paragraph.text)
        # Retain heading info if needed for context (headings are critical)
        if paragraph.style.name.startswith('Heading'):
            text.append(f"HEADING: {paragraph.text}")
    return text


# Function to load text from PowerPoint (PPTX)
def load_ppt(file_path):
    pythoncom.CoInitialize()  # Initialize COM for PowerPoint reading
    presentation = Presentation(file_path)
    text = []
    for slide in presentation.slides:
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_text.append(shape.text)
        if slide_text:
            text.append("\n".join(slide_text))  # Keep text organized by slide
    return text


# Function to load text from plain Text files (TXT)
def load_txt(file_path):
    with open(file_path, 'r', encoding='utf-8') as f:
        return f.readlines()


# Enhanced text preprocessing function to retain meaningful context
def preprocess_text(text):
    # Ensure that text is a string
    if not isinstance(text, str):
        text = str(text)  # Convert non-string types to string

    # Clean text by removing unwanted characters but retain structured data (headings, sections)
    cleaned_text = re.sub(r'[^A-Za-z0-9\s\.\,\;\:\-\'\(\)]', '', text)  # Keep punctuation for structure
    cleaned_text = cleaned_text.lower()  # Convert to lowercase
    stop_words = set(stopwords.words('english'))  # List of stopwords
    filtered_words = [word for word in cleaned_text.split() if word not in stop_words]
    return ' '.join(filtered_words)


# Function to chunk documents based on sections or headings (optimized for academic content)
def chunk_document(doc, text_splitter):
    """Split documents into meaningful chunks based on context."""
    return text_splitter.split_documents([doc])


# Enhanced setup_vectorstore function with chunk size adjustments and better chunk overlap
def setup_vectorstore(documents):
    embeddings = HuggingFaceEmbeddings(model_name="allenai/scibert_scivocab_uncased")  # SciBERT for better academic content understanding
    text_splitter = CharacterTextSplitter(
        separator="\n",  # Split documents by newlines for structured content
        chunk_size=1000,  # Increase chunk size to ensure that chunks are complete
        chunk_overlap=200  # Larger overlap to retain context across chunks
    )

    # Parallel processing of document chunks
    with concurrent.futures.ThreadPoolExecutor() as executor:
        doc_chunks = list(executor.map(lambda doc: chunk_document(doc, text_splitter), documents))

    # Flatten the list of chunks
    flattened_chunks = [chunk for sublist in doc_chunks for chunk in sublist]
    vectorstore = FAISS.from_documents(flattened_chunks, embeddings)
    return vectorstore


# Function to create the conversation chain with better parameters for complete answers
def create_chain(vectorstore, creativity_level, num_chunks, response_type):
    llm = ChatGroq(
        model="llama-3.3-70b-versatile",  # Ensure you're using a detailed model
        temperature=creativity_level
    )

    # Increase k value to allow the model to access more relevant chunks
    retriever = vectorstore.as_retriever(search_type="similarity", search_kwargs={"k": num_chunks * 2})

    # Use 'stuff' instead of 'refine' to generate more concise answers without extra explanation
    chain = ConversationalRetrievalChain.from_llm(
        llm=llm,
        retriever=retriever,
        chain_type="stuff",  # Use stuff for more direct answers without refinement process
        memory=None,
        verbose=False  # Disable verbose logging
    )
    return chain


# Function to handle voice input
def handle_voice_query():
    recognizer = sr.Recognizer()
    with sr.Microphone() as source:
        st.info("Listening for your question... Speak now.")
        try:
            audio = recognizer.listen(source, timeout=5)
            query = recognizer.recognize_google(audio)
            return query
        except sr.UnknownValueError:
            st.error("Sorry, I couldn't understand what you said. Please try again.")
        except sr.RequestError as e:
            st.error(f"Could not request results; {e}")
        except Exception as e:
            st.error(f"An error occurred: {e}")
    return None


# Set up Streamlit page
st.set_page_config(
    page_title="Chat with Doc",
    page_icon="📄",
    layout="centered"
)

st.markdown("""  
<div style="text-align: center;">
    <h1>Welcome to SmartDOC 📄</h1>
    <h5><i>Transform Your Documents into Intelligent Conversations 💬</i></h5>
</div>
""", unsafe_allow_html=True)

with st.expander("**About this app**"):
    st.write(""" 
        - Leverages the LLAMA 3.3 70b-versatile model to quickly transform documents (PDF, Word, PPT, Text) into interactive conversations.
        - Extract insights and get precise answers based on the document’s content.
        - Save time with context-specific responses without reading everything.
        - Provides fast and accurate answers in real-time.
    """)

# Initialize chat history in Streamlit session state
if "chat_history" not in st.session_state:
    st.session_state.chat_history = []

# Sidebar: Add sliders for creativity level and number of retrieved chunks
with st.sidebar:
    st.header("Response Parameters")
    creativity_level = st.slider("**Creativity Level**", 0.00, 1.0, 0.2, 0.01)
    num_chunks = st.slider("**Number of Retrieved Chunks**", 1, 10, 5)
    response_type = st.radio("**Response Type**", ["Brief", "Elaborate"])

# Layout: File upload section
uploaded_files = st.file_uploader(label="Upload your documents", type=["pdf", "docx", "pptx", "txt"],
                                  accept_multiple_files=True)

# Initialize user_input as None
user_input = None

# Document processing and question-answering logic after file upload
if uploaded_files:
    # Save files to disk
    file_paths = []
    for uploaded_file in uploaded_files:
        file_path = os.path.join(working_dir, uploaded_file.name)
        file_paths.append(file_path)
        with open(file_path, "wb") as f:
            f.write(uploaded_file.getbuffer())

    # Show loading spinner while the documents are being processed
    with st.spinner("Processing documents... This may take a while."):

        # Load and process all documents efficiently with parallelism
        all_documents = []
        with concurrent.futures.ThreadPoolExecutor() as executor:
            all_documents = list(executor.map(load_document, file_paths))

        # Flatten the list of documents
        all_documents = [doc for sublist in all_documents for doc in sublist]

        # Preprocess documents for better extraction
        all_documents = [CustomDocument(preprocess_text(doc.page_content)) for doc in all_documents]

        # Setup vector store and conversation chain
        if "vectorstore" not in st.session_state:
            st.session_state.vectorstore = setup_vectorstore(all_documents)

        if "conversation_chain" not in st.session_state:
            st.session_state.conversation_chain = create_chain(st.session_state.vectorstore, creativity_level, num_chunks, response_type)

        # Log the processing time
        end_time = time.time()
        start_time = time.time()  # Define start time
        st.success(f"Documents processed in {end_time - start_time:.2f} seconds")

    # Show the input fields for asking questions and voice input after documents are uploaded
    user_input = st.text_input("Enter your questions:")

    # Voice input button
    if st.button("Ask via Voice"):
        voice_query = handle_voice_query()
        if voice_query:
            user_input = voice_query
            st.info(f"Question: {user_input}")
            # Display the notification below the question
            st.success(f"Answer for your question '{user_input}' has been generated.")

# Display chat history
for message in st.session_state.chat_history:
    with st.chat_message(message["role"]):
        st.markdown(message["content"])

# Handling user inputs (queries)
if user_input:
    # Display the user input as a chat message
    with st.chat_message("user"):
        st.markdown(user_input)

    # Prepare chat history for conversation chain
    formatted_chat_history = [(msg["role"], msg["content"]) for msg in st.session_state.chat_history]

    # Prepare query with formatted chat history
    query_with_history = {
        "question": user_input,
        "chat_history": formatted_chat_history  # Pass the correctly formatted history here
    }

    try:
        # Process the query with the conversation chain
        response = st.session_state.conversation_chain(query_with_history)
        assistant_response = response["answer"]

        # Modify response type based on user selection
        if response_type == "Brief":
            assistant_response = assistant_response[:500]  # Limit to first 500 characters

        # Provide the final cleaned answer
        with st.chat_message("assistant"):
            st.markdown(assistant_response)
            # Append both question and answer to the chat history sequentially
            st.session_state.chat_history.append({"role": "user", "content": user_input})
            st.session_state.chat_history.append({"role": "assistant", "content": assistant_response})

    except Exception as e:
        st.error(f"Error during processing: {e}")
